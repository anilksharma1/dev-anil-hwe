"""Per-format, read-only text extraction.

Every extractor returns (text, meta). No extractor executes document logic
(macros, formulas, scripts); none writes to disk. Failures raise and are
isolated by the caller. `meta` may carry: text_extractable, is_structured,
page_or_sheet_count, attachment_count, structured_entity_rows,
estimate_truncated, status, detail.
"""
from __future__ import annotations

import os
import re
import csv as _csv
import signal
import time
import zipfile
import contextlib
import logging as _logging

# pypdf logs a WARNING per undecodable XObject; a scanned PDF emits thousands, and across
# many workers that stderr/logging flood is real overhead (and noise). Silence it -- we
# already handle unreadable content by routing to OCR.
_logging.getLogger("pypdf").setLevel(_logging.ERROR)
_log = _logging.getLogger(__name__)

from .detection import CompiledRules, row_has_identifier, is_roster_name_line


class _NullStderr:
    def write(self, *_a):
        pass

    def flush(self):
        pass


@contextlib.contextmanager
def _hush_stderr():
    """Swallow stderr for the duration of a block. Some pypdf builds print benign
    'Impossible to decode XFormObject' chatter straight to stderr regardless of
    logging level; this guarantees it never reaches the console. Real extraction
    errors are caught by the caller and reported via status/detail, not stderr."""
    with contextlib.redirect_stderr(_NullStderr()):
        yield


def quiet_noisy_libraries() -> None:
    """Silence benign third-party chatter that floods the console but means nothing:
    pypdf's 'Impossible to decode XFormObject' (an embedded image/form it skips while
    still reading the page), pdfminer chatter, and openpyxl's 'Data Validation extension
    is not supported' UserWarning. None affect extraction or detection. Called per worker
    (and in the parent)."""
    import logging
    import warnings
    for name in ("pypdf", "pypdf._page", "PyPDF2",
                 "pdfminer", "pdfplumber", "pdfminer.pdfinterp", "pdfminer.pdfpage"):
        logging.getLogger(name).setLevel(logging.ERROR)
    warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")


# Optional dependencies (probed once).
try:
    import PIL  # noqa: F401  -- pypdf needs Pillow to decode embedded images
    HAVE_PILLOW = True
except Exception:
    HAVE_PILLOW = False

try:
    import pypdf
    HAVE_PYPDF = True
except Exception:
    HAVE_PYPDF = False
try:
    import docx
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False
try:
    import openpyxl
    HAVE_OPENPYXL = True
except Exception:
    HAVE_OPENPYXL = False
try:
    from pptx import Presentation
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False
try:
    import extract_msg
    HAVE_MSG = True
except Exception:
    HAVE_MSG = False
try:
    import xlrd
    HAVE_XLRD = True
except Exception:
    HAVE_XLRD = False
try:
    from striprtf.striprtf import rtf_to_text
    HAVE_STRIPRTF = True
except Exception:
    HAVE_STRIPRTF = False


class TimeoutError_(Exception):
    pass


class time_limit:
    """Per-file wall-clock guard (Unix; no-op elsewhere)."""

    def __init__(self, seconds: int):
        self.seconds = seconds
        self.enabled = seconds > 0 and hasattr(signal, "SIGALRM")

    def __enter__(self):
        if self.enabled:
            signal.signal(signal.SIGALRM, self._raise)
            signal.alarm(self.seconds)
        return self

    def __exit__(self, *exc):
        if self.enabled:
            signal.alarm(0)
        return False

    @staticmethod
    def _raise(signum, frame):
        raise TimeoutError_()


def _safe_zip_text(path, name_filter, cfg, char_cap):
    """Concatenate text from selected zip entries, with a zip-bomb guard."""
    chunks, total = [], 0
    with zipfile.ZipFile(path) as zf:
        for info in zf.infolist():
            if info.file_size and info.compress_size:
                if info.file_size / max(info.compress_size, 1) > cfg.zip_ratio_limit:
                    raise ValueError("decompression ratio guard tripped")
            if not name_filter(info.filename):
                continue
            with zf.open(info) as fh:
                data = fh.read(max(char_cap - total, 0) + 1)
            chunks.append(data.decode("utf-8", "ignore"))
            total += len(data)
            if total >= char_cap:
                break
    return re.sub(r"<[^>]+>", " ", " ".join(chunks))


def _strip_html(raw: str) -> str:
    raw = re.sub(r"(?is)<(script|style)\b.*?</\1>", " ", raw)
    return re.sub(r"<[^>]+>", " ", raw)


# --------------------------------------------------------------------------- #
# Text / delimited
# --------------------------------------------------------------------------- #
def x_text(path, cfg, rules):
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        return fh.read(cfg.max_scan_chars), {"text_extractable": "text"}


def x_delimited(path, cfg, rules):
    """CSV/TSV: treat as structured; count rows that represent a person."""
    delim = "\t" if path.lower().endswith((".tsv", ".tab")) else ","
    entity_rows, name_rows, scanned, truncated, sample = 0, 0, 0, False, []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as fh:
        reader = _csv.reader(fh, delimiter=delim)
        for row in reader:
            scanned += 1
            text = " ".join(row)
            if scanned <= 50:
                sample.append(text)
            if row_has_identifier(text, rules):
                entity_rows += 1
            elif is_roster_name_line(text):
                name_rows += 1
            if scanned >= cfg.max_scan_rows:
                truncated = True
                break
    return "\n".join(sample), {
        "text_extractable": "text", "is_structured": True,
        "structured_entity_rows": entity_rows + name_rows, "structured_total_rows": scanned,
        "estimate_truncated": truncated,
    }


# --------------------------------------------------------------------------- #
# OOXML (zip-based)
# --------------------------------------------------------------------------- #
def x_docx(path, cfg, rules):
    if HAVE_DOCX:
        try:
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)[: cfg.max_scan_chars], \
                {"text_extractable": "text"}
        except Exception as exc:
            _log.debug("docx native parse failed for %s, falling back to zip: %s", path, exc, exc_info=True)
    text = _safe_zip_text(path, lambda n: n.startswith("word/") and n.endswith(".xml"),
                          cfg, cfg.max_scan_chars)
    return text, {"text_extractable": "text"}


def x_pptx(path, cfg, rules):
    if HAVE_PPTX:
        try:
            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)[: cfg.max_scan_chars], \
                {"text_extractable": "text", "page_or_sheet_count": len(prs.slides._sldIdLst)}
        except Exception as exc:
            _log.debug("pptx native parse failed for %s, falling back to zip: %s", path, exc, exc_info=True)
    text = _safe_zip_text(path, lambda n: n.startswith("ppt/slides/") and n.endswith(".xml"),
                          cfg, cfg.max_scan_chars)
    return text, {"text_extractable": "text"}


def x_xlsx(path, cfg, rules):
    if not HAVE_OPENPYXL:
        return "", {"text_extractable": "unknown", "status": "no_parser", "is_structured": True}
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    entity_rows, name_rows, scanned, truncated, sample = 0, 0, 0, False, []
    sheets = len(wb.worksheets)
    for ws in wb.worksheets:
        # Some workbooks (notably payroll/HR registers) declare a broken sheet dimension
        # -- often flagged by openpyxl's "Print area cannot be set to Defined name" warning
        # -- and in read_only mode iter_rows then trusts that dimension and returns almost
        # nothing (observed: 3 rows from a 3,012-row, 600-person register). Resetting the
        # dimension forces openpyxl to read every row that is actually present.
        try:
            ws.reset_dimensions()
        except Exception:
            pass
        for row in ws.iter_rows(values_only=True):
            scanned += 1
            text = " ".join("" if c is None else str(c) for c in row)
            if scanned <= 50:
                sample.append(text)
            if row_has_identifier(text, rules):
                entity_rows += 1
            elif is_roster_name_line(text):        # name-only person row (register layout)
                name_rows += 1
            if scanned >= cfg.max_scan_rows:
                truncated = True
                break
        if truncated:
            break
    wb.close()
    # Subject count = identifier rows PLUS name-only roster rows (the elif above means a
    # row is counted once). Recovers registers where each person's row carries a name but
    # no other identifier on that same row.
    entity_estimate = entity_rows + name_rows
    return "\n".join(sample), {
        "text_extractable": "text", "is_structured": True, "page_or_sheet_count": sheets,
        "structured_entity_rows": entity_estimate, "structured_total_rows": scanned,
        "estimate_truncated": truncated,
    }


def x_odf_text(path, cfg, rules):
    """OpenDocument text/presentation: pull content.xml as free text."""
    text = _safe_zip_text(path, lambda n: n == "content.xml", cfg, cfg.max_scan_chars)
    return text, {"text_extractable": "text"}


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #
def _pdf_field_values(reader) -> str:
    """Typed values in a fillable PDF live in AcroForm fields, not the page text.
    Pull them in as 'field name: value' so a form's name/SSN/email are visible to
    detection (and the field name supplies the label for value-near-label rules)."""
    try:
        fields = reader.get_fields()
    except Exception:
        return ""
    if not fields:
        return ""
    out = []
    for name, f in fields.items():
        try:
            v = f.get("/V")
        except Exception:
            v = None
        if v is None:
            continue
        s = str(v).strip()
        if s:
            out.append(f"{name}: {s}" if name else s)
    return "\n".join(out)


def _pdf_image_based(reader, pages: int) -> bool:
    """True when most sampled pages are dominated by a LARGE (roughly full-page) image
    -- a scan or image overlay whose content OCR must read. A small logo or letterhead
    graphic does NOT count, so ordinary text documents are never sent to OCR by mistake."""
    sample = min(pages, 6)
    if sample <= 0:
        return False
    big_image_pages = 0
    for page in reader.pages[:sample]:
        try:
            res = page.get("/Resources")
            xobj = res.get("/XObject").get_object() if res and res.get("/XObject") else None
            if not xobj:
                continue
            for k in xobj:
                o = xobj[k].get_object()
                if o.get("/Subtype") == "/Image":
                    w = float(o.get("/Width", 0) or 0)
                    h = float(o.get("/Height", 0) or 0)
                    if w >= 550 and h >= 550:   # full-page scan, not a logo/banner
                        big_image_pages += 1
                        break
        except Exception:
            continue
    return big_image_pages >= (sample + 1) // 2


def _pdf_extract_content_images(reader):
    """Extract (1-indexed page number, JPEG bytes) pairs for content-sized embedded
    images in a text PDF.

    Skips images outside Azure DI's accepted range (50–10000 px per side), logos/
    decorations (w < 400), and full-page scans (area >= 550*550, already handled as
    image_only). Deduplicates by XObject name; records the first page each image
    appears on. Returns [] on any failure.

    3.0.0 -- PILLOW. `img.image` below requires Pillow. Without it EVERY image raises
    and is swallowed by the inner `except Exception: continue`, so this returns [] and
    apply_image_ocr silently becomes a no-op. That is not hypothetical: Daniel's run
    logged 0 embedded-image OCR calls against Anna's 24,223 on the identical corpus,
    because Pillow was absent and nothing said so. Pillow is now declared in
    requirements.txt, reported by optional_dependency_report(), and the decode-failure
    count is returned so a broken install shows up in the inventory instead of looking
    like a cheap Azure bill.

    Returns (content_images, stats) where stats = {"seen": int, "decode_failed": int}.
    """
    import io
    _DI_MIN, _DI_MAX = 50, 10_000
    content_images = []   # list of (page_1indexed: int, img_bytes: bytes)
    stats = {"seen": 0, "decode_failed": 0}
    seen = set()
    for page_idx, page in enumerate(reader.pages):
        try:
            for img in page.images:
                name = getattr(img, "name", None) or ""
                if name in seen:
                    continue
                stats["seen"] += 1
                try:
                    pil = img.image
                    if pil is None:
                        continue
                    w, h = pil.size
                    if not (_DI_MIN <= w <= _DI_MAX and _DI_MIN <= h <= _DI_MAX
                            and w >= 400 and w * h < 550 * 550):
                        continue
                    seen.add(name)
                    buf = io.BytesIO()
                    pil.convert("RGB").save(buf, format="JPEG")
                    content_images.append((page_idx + 1, buf.getvalue()))
                except Exception as exc:
                    # Almost always a missing/broken Pillow. Counted, not silent.
                    stats["decode_failed"] += 1
                    _log.debug("pdf embedded image decode failed (page %d): %s", page_idx + 1, exc, exc_info=True)
                    continue
        except Exception as exc:
            _log.debug("pdf image iteration failed for page %d of %s: %s", page_idx + 1, path, exc, exc_info=True)
            continue
    return content_images, stats


def _call_with_timeout(fn, seconds):
    """Run fn() with a hard wall-clock cap via a daemon thread (SIGALRM is Unix-only and
    inert on the Windows workers). On timeout the call is abandoned and TimeoutError is
    raised, so a pypdf open that hangs on a malformed PDF returns fast instead of letting
    the parent watchdog kill the whole worker. The abandoned thread is bounded by the
    parent watchdog / worker recycle as a backstop."""
    import threading
    box = {}

    def run():
        try:
            box["ok"] = fn()
        except BaseException as exc:  # noqa: BLE001
            box["err"] = exc

    t = threading.Thread(target=run, daemon=True)
    t.start()
    t.join(seconds)
    if t.is_alive():
        raise TimeoutError(f"pdf open exceeded {seconds:.0f}s")
    if "err" in box:
        raise box["err"]
    return box.get("ok")


_RAW_PAGE = re.compile(rb"/Type\s*/Page[^s]")


def _pdf_page_count_raw(path):
    """Approximate page count straight from the file bytes, WITHOUT pypdf -- so we still
    have a total for extrapolation on a PDF that pypdf can't open. Returns 0 if unknown."""
    try:
        with open(path, "rb") as fh:
            data = fh.read()
        n = len(_RAW_PAGE.findall(data))
        return n if n > 0 else 0
    except Exception:
        return 0


def x_pdf(path, cfg, rules):
    if not HAVE_PYPDF:
        return "", {"text_extractable": "unknown", "status": "no_parser"}
    timing = bool(getattr(cfg, "phase_timing", False))
    # Cap how many pages we pull text from. A genuine text PDF hits the max_scan_chars
    # break long before this; the cap only bounds pathological/huge scans where pypdf's
    # per-page work would otherwise run unbounded and blow the parent watchdog. The kind
    # (text vs image_only) classification below is unchanged: a capped scan still yields
    # ~no text, so it is still detected as image_only. 0 disables the page cap.
    page_cap = int(getattr(cfg, "pdf_max_extract_pages", 0) or 0)
    t0 = time.monotonic()
    # Bound the pypdf OPEN: ~1% of scanned PDFs hang pypdf here and never reach OCR (the
    # worker gets watchdog-killed). On timeout/failure, return a 'pdf_unreadable' signal
    # fast -- apply_ocr then sends the file straight to Azure DI (no pypdf), page-capped,
    # and a BDE count is recovered (promote-only). page_or_sheet_count comes from a
    # pypdf-free byte scan so extrapolation still has a total.
    open_timeout = float(getattr(cfg, "pdf_open_timeout", 30) or 30)

    def _open():
        r = pypdf.PdfReader(path)
        return r, len(r.pages)

    try:
        with _hush_stderr():
            reader, pages = _call_with_timeout(_open, open_timeout)
    except Exception as exc:  # noqa: BLE001 -- hang, corruption, or parse failure
        _log.warning("pdf open failed for %s: %s", path, exc, exc_info=True)
        return "", {"text_extractable": "pdf_unreadable",
                    "page_or_sheet_count": _pdf_page_count_raw(path),
                    "detail": f"pdf_open_failed:{type(exc).__name__}"}
    with _hush_stderr():
        t_open = time.monotonic()
        # For a likely scan, stop the extract_text sweep after the first few pages: a scan
        # yields ~no text no matter how many pages we read, and the full sweep over 100s of
        # pages (x many workers on a network share) is exactly what blows the watchdog. We
        # only cut out early once the file is confirmed image-based, so genuine text PDFs
        # (which accumulate text fast) read to completion unchanged.
        probe_pages = int(getattr(cfg, "pdf_scan_probe_pages", 0) or 0)
        image_based = None
        parts, total, read_pages = [], 0, 0
        for page in reader.pages:
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            parts.append(t)
            total += len(t)
            read_pages += 1
            if total >= cfg.max_scan_chars:
                break
            if page_cap and read_pages >= page_cap:
                break
            # sparse after the probe window -> check once if it's a scan; if so, stop early.
            if probe_pages and read_pages >= probe_pages and total < 20 * read_pages:
                if image_based is None:
                    image_based = _pdf_image_based(reader, pages)
                if image_based:
                    break
        text = "\n".join(parts)
        t_text = time.monotonic()
        field_text = _pdf_field_values(reader)        # fillable-form values, if any
        if image_based is None:
            image_based = _pdf_image_based(reader, pages)
        t_img = time.monotonic()
    if field_text:
        text = (text + "\n" + field_text)[: cfg.max_scan_chars]

    # Decide whether OCR is needed. A genuine text PDF has substantial text per page;
    # a scan/overlay form has a thin or label-only layer with none of the filled values.
    body_len = len(text.strip())
    try:
        size_bytes = os.path.getsize(path)
    except OSError:
        size_bytes = 0
    low_density = pages > 0 and body_len < max(20, 50 * pages)
    # Image-dominated: a lot of bytes (an embedded scan) but very little extracted
    # text -- a scanned page with a junk/partial text layer that the thin-text and
    # full-page-image checks can both miss. Genuine text PDFs (lots of text) are
    # left alone by the body_len < 1500 guard, so this won't re-OCR ordinary docs.
    image_dominated = body_len < 1500 and size_bytes / max(body_len, 1) > 500
    kind = "image_only" if (low_density or image_dominated or image_based) else "text"
    if kind == "text":
        # Tolerate a patched/older _pdf_extract_content_images that returns a bare list:
        # several tests monkeypatch this symbol, and the stats tuple is 3.0.0-only.
        _res = _pdf_extract_content_images(reader)
        if isinstance(_res, tuple):
            content_images, _img_stats = _res
        else:
            content_images, _img_stats = _res, {"seen": len(_res or []), "decode_failed": 0}
    else:
        content_images, _img_stats = [], {"seen": 0, "decode_failed": 0}
    meta = {"text_extractable": kind, "page_or_sheet_count": pages,
            "content_images": content_images,
            "pdf_images_seen": _img_stats["seen"],
            "pdf_images_decode_failed": _img_stats["decode_failed"]}
    notes = []
    if page_cap and read_pages >= page_cap and read_pages < pages:
        notes.append(f"pdf_text_capped:{read_pages}/{pages}")
    if timing:
        notes.append(f"pdf_open_ms:{int((t_open - t0) * 1000)} "
                     f"pdf_text_ms:{int((t_text - t_open) * 1000)} "
                     f"pdf_imgdetect_ms:{int((t_img - t_text) * 1000)}")
    if notes:
        meta["detail"] = " ".join(notes)
    return text, meta


# --------------------------------------------------------------------------- #
# Email
# --------------------------------------------------------------------------- #
def x_msg(path, cfg, rules):
    if not HAVE_MSG:
        return "", {"text_extractable": "unknown", "status": "no_parser"}
    m = extract_msg.Message(path)
    subject = m.subject or ""
    body = m.body or ""
    # Also read the HTML body so content that exists only in HTML -- notably
    # emails/phones that appear as links in signatures and forwarded chains --
    # is captured, not just the plain-text body.
    html = m.htmlBody or b""
    if isinstance(html, bytes):
        html = html.decode("utf-8", "ignore")
    html_text = ""
    if html:
        links = " ".join(re.findall(r"(?:mailto:|tel:)([^\"'>\s]+)", html, re.I))
        html_text = _strip_html(html) + " " + links
    n_att = len(m.attachments) if m.attachments is not None else 0
    m.close()
    combined = f"{subject}\n{body}\n{html_text}"[: cfg.max_scan_chars]
    return combined, {"text_extractable": "text", "attachment_count": n_att}


def x_eml(path, cfg, rules):
    import email
    from email import policy
    with open(path, "rb") as fh:
        msg = email.message_from_binary_file(fh, policy=policy.default)
    parts, n_att = [], 0
    parts.append(msg.get("subject", "") or "")
    if msg.is_multipart():
        for part in msg.walk():
            disp = (part.get_content_disposition() or "")
            ctype = part.get_content_type()
            if disp == "attachment" or part.get_filename():
                n_att += 1
                continue
            if ctype == "text/plain":
                parts.append(part.get_content())
            elif ctype == "text/html":
                parts.append(_strip_html(part.get_content()))
    else:
        body = msg.get_content()
        parts.append(_strip_html(body) if msg.get_content_type() == "text/html" else body)
    text = "\n".join(str(p) for p in parts)[: cfg.max_scan_chars]
    return text, {"text_extractable": "text", "attachment_count": n_att}


# --------------------------------------------------------------------------- #
# Markup
# --------------------------------------------------------------------------- #
def x_html(path, cfg, rules):
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        return _strip_html(fh.read(cfg.max_scan_chars * 2))[: cfg.max_scan_chars], \
            {"text_extractable": "text"}


def x_rtf(path, cfg, rules):
    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
        raw = fh.read(cfg.max_scan_chars * 2)
    if HAVE_STRIPRTF:
        text = rtf_to_text(raw, errors="ignore")
    else:  # stdlib fallback: drop groups/control words
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", raw)
        text = re.sub(r"[{}]", " ", text)
    return text[: cfg.max_scan_chars], {"text_extractable": "text"}


# --------------------------------------------------------------------------- #
# Legacy binary office (need converters)
# --------------------------------------------------------------------------- #
def x_doc(path, cfg, rules):
    import shutil
    import subprocess
    tool = shutil.which("antiword") or shutil.which("catdoc")
    if not tool:
        return "", {"text_extractable": "needs_conversion", "status": "no_parser"}
    try:
        out = subprocess.run([tool, path], capture_output=True, timeout=cfg.timeout_s,
                             check=False)
        return out.stdout.decode("utf-8", "ignore")[: cfg.max_scan_chars], \
            {"text_extractable": "text"}
    except Exception as exc:
        _log.debug("antiword/catdoc failed for %s: %s", path, exc, exc_info=True)
        return "", {"text_extractable": "needs_conversion", "status": "no_parser"}


def x_xls(path, cfg, rules):
    if not HAVE_XLRD:
        return "", {"text_extractable": "needs_conversion", "status": "no_parser",
                    "is_structured": True}
    book = xlrd.open_workbook(path)
    entity_rows, name_rows, scanned, truncated, sample = 0, 0, 0, False, []
    for sh in book.sheets():
        for r in range(sh.nrows):
            scanned += 1
            text = " ".join(str(v) for v in sh.row_values(r))
            if scanned <= 50:
                sample.append(text)
            if row_has_identifier(text, rules):
                entity_rows += 1
            elif is_roster_name_line(text):
                name_rows += 1
            if scanned >= cfg.max_scan_rows:
                truncated = True
                break
        if truncated:
            break
    return "\n".join(sample), {
        "text_extractable": "text", "is_structured": True, "page_or_sheet_count": book.nsheets,
        "structured_entity_rows": entity_rows + name_rows, "structured_total_rows": scanned,
        "estimate_truncated": truncated,
    }


def x_needs_conversion(path, cfg, rules):
    """Formats we can recognize but not parse here (legacy ppt, etc.)."""
    return "", {"text_extractable": "needs_conversion", "status": "no_parser"}


# --------------------------------------------------------------------------- #
# Images (flagged for OCR, never read here) and containers (flagged to expand)
# --------------------------------------------------------------------------- #
def x_image(path, cfg, rules):
    return "", {"text_extractable": "image_only"}


def x_container(path, cfg, rules):
    """Mail stores / archives: flagged for expansion, not extracted in triage."""
    return "", {"text_extractable": "none", "status": "container", "detail": "expand_first"}


# --------------------------------------------------------------------------- #
# Registry
# --------------------------------------------------------------------------- #
EXTRACTORS = {
    ".txt": x_text, ".log": x_text, ".text": x_text,
    ".csv": x_delimited, ".tsv": x_delimited, ".tab": x_delimited,
    ".docx": x_docx, ".docm": x_docx,
    ".pptx": x_pptx, ".pptm": x_pptx,
    ".xlsx": x_xlsx, ".xlsm": x_xlsx,
    ".odt": x_odf_text, ".odp": x_odf_text,
    ".pdf": x_pdf,
    ".msg": x_msg, ".eml": x_eml,
    ".html": x_html, ".htm": x_html, ".xml": x_html,
    ".rtf": x_rtf,
    ".doc": x_doc, ".xls": x_xls, ".ppt": x_needs_conversion, ".ods": x_needs_conversion,
    ".png": x_image, ".jpg": x_image, ".jpeg": x_image, ".tif": x_image,
    ".tiff": x_image, ".bmp": x_image, ".gif": x_image,
    ".zip": x_container, ".pst": x_container, ".ost": x_container,
    ".nsf": x_container, ".mbox": x_container, ".7z": x_container, ".rar": x_container,
}


def get_extractor(ext: str):
    return EXTRACTORS.get(ext.lower())


def optional_dependency_report() -> list[str]:
    missing = []
    for name, ok in (("pypdf", HAVE_PYPDF), ("python-docx", HAVE_DOCX),
                     ("openpyxl", HAVE_OPENPYXL), ("python-pptx", HAVE_PPTX),
                     ("extract-msg", HAVE_MSG), ("xlrd", HAVE_XLRD),
                     ("striprtf", HAVE_STRIPRTF), ("Pillow", HAVE_PILLOW)):
        if not ok:
            missing.append(name)
    return missing